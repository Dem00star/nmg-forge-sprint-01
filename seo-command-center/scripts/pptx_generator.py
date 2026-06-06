import json
from pptx import Presentation

def generate_pptx(json_string):
    data = json.loads(json_string)
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = data.get('site', 'SEO Audit')
    subtitle.text = "Automated SEO Audit"

    # Executive Summary Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Overview"
    body = slide.placeholders[1]
    summary_data = data.get('summary', {})
    body.text = f"URLs Crawled: {data.get('urls_crawled', 'N/A')}\nTotal Issues: {summary_data.get('total_issues', 'N/A')}"

    # Issues Slides
    issues = data.get('issues', [])
    for issue in issues:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = issue.get('type', 'Issue')
        body = slide.placeholders[1]
        body.text = f"Severity: {issue.get('severity', 'N/A')}\nAffected URLs: {issue.get('count', 'N/A')}"

    prs.save('final_report.pptx')
